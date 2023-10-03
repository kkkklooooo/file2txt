
import PyPDF2
import docx
import pptx
import os
import time
import shutil
import sys
import argparse


def sss(path):
    try:
        c = 0
        text = []
        reader = PyPDF2.PdfReader(path)
        for index, page in enumerate(reader.pages):
            c += 1
            text.append(page.extract_text())  # 遍历所有页面
            # print(page.extract_text())
            if c > 1:
                pass
                # break  # 仅load第一页
        f = ''.join(text)
        # print(f)
        if f == "":
            f = "None"
        return f
    except:
        return "fuck"


def readppt(path):
    ppt = pptx.Presentation(path)
    results = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    part = []
                    for run in paragraph.runs:
                        part.append(run.text)
                    results.append(''.join(part))
    return "".join(results)


def readdocx(path):
    try:
        doc = docx.Document(path)
        text = []
        for par in doc.paragraphs:
            if par.text:
                text.append(par.text)
        return ''.join(text)
    except:
        return "fuck"


def chose_func(extension, full_path):
    if extension == ".pdf":
        return sss(full_path)
    elif extension == ".pptx":
        return readppt(full_path)
    elif extension == ".docx":
        return readdocx(full_path)
    else:
        return "None"


def save(path, data):
    with open(path, "w", encoding='utf-8') as f:
        f.write(data)


parser = argparse.ArgumentParser(description="FileToTxt")
parser.add_argument("-t", "--type", type=int,
                    help="指定转换file还是path(0:file,1:path)")
# if parser.parse_args().type == 0:
parser.add_argument("-i", "--input", type=str, help="输入文件或路径")
parser.add_argument("-o", "--output", type=str, help="输出文件路径")
parser.add_argument("-e", "--extension", type=str, help="指定文件类型(仅当type=1时有效)")
args = parser.parse_args()

file = args.input
if args.type == 0:
    opt_filename, extension = os.path.splitext(file)
    data = chose_func(extension, file)
    save(opt_filename+".txt", data)
    print(f"finish {opt_filename+extension}\n 字符数{len(data)}")
else:
    have_e = args.extension != None
    if os.path.isdir(file):
        dir = file
    elif os.path.isfile(file):
        dir = os.path.dirname(file)
    else:
        print('Input Error')
    files = os.listdir(dir)
    for i in files:
        opt_filename, extension = os.path.splitext(i)
        if have_e:
            if args.extension in extension:
                data = chose_func(extension, dir+"\\"+i)
                save(dir+"\\"+opt_filename+".txt", data)
                print(f"finish {i}\n 字符数{len(data)}")
        else:
            print("去他妈的")
